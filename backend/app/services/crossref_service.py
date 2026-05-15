import os
import csv
import io
import re
from pathlib import Path
from typing import Any
from app.core.config import settings


from app.services.llm_service import extract_fields


class CrossrefService:
    SUPPORTED_EXTENSIONS = {".pdf", ".csv", ".pptx", ".ppt", ".docx", ".doc"}

    def __init__(self):
        self.upload_dir = Path(settings.upload_dir) / "crossref"
        self.upload_dir.mkdir(parents=True, exist_ok=True)

    async def semantic_match(self, query_data: dict, candidates: list[dict], output_columns: list[str]) -> dict:
        """Usa Gemini 2.0 Flash para encontrar el mejor match semántico."""
        prompt = f"""
        Eres un experto en cruce de datos. Debes encontrar cuál de los 'Candidatos' coincide mejor con el 'Origen'.
        
        Origen: {json.dumps(query_data)}
        
        Candidatos:
        {json.dumps(candidates[:10], indent=2)}
        
        Si encuentras un match claro, devuelve un JSON con los campos solicitados: {output_columns}.
        Si no hay match, devuelve todos los campos como "NO ENCONTRADO".
        """
        
        schema = {
            "type": "object",
            "properties": {col: {"type": "string"} for col in output_columns}
        }
        
        try:
            result = await extract_fields(prompt, schema, model=settings.gemini_model_crossref)
            return result
        except:
            return {col: "NO ENCONTRADO" for col in output_columns}

    def save_file(self, file_content: bytes, filename: str) -> str:
        dest = self.upload_dir / filename
        with open(dest, "wb") as f:
            f.write(file_content)
        return str(dest)

    def parse_file(self, file_path: str) -> tuple[list[str], list[dict]]:
        ext = Path(file_path).suffix.lower()
        if ext == ".csv":
            return self._parse_csv(file_path)
        elif ext == ".pdf":
            return self._parse_pdf(file_path)
        elif ext in (".pptx", ".ppt"):
            return self._parse_ppt(file_path)
        elif ext in (".docx", ".doc"):
            return self._parse_docx(file_path)
        else:
            raise ValueError(f"Formato no soportado: {ext}")

    def _parse_csv(self, file_path: str) -> tuple[list[str], list[dict]]:
        import pandas as pd
        encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]
        df = None
        for enc in encodings:
            try:
                df = pd.read_csv(file_path, dtype=str, encoding=enc, sep=None, engine="python")
                break
            except (UnicodeDecodeError, pd.errors.ParserError):
                continue
        if df is None:
            df = pd.read_csv(file_path, dtype=str, encoding="latin-1", on_bad_lines="skip")
        df = df.fillna("")
        columns = list(df.columns)
        rows = df.to_dict(orient="records")
        return columns, rows

    def _parse_pdf(self, file_path: str) -> tuple[list[str], list[dict]]:
        import pdfplumber
        rows: list[dict] = []
        columns: list[str] = []

        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                table = page.extract_table()
                if table:
                    if not columns:
                        columns = [c.strip() if c else f"col_{i}" for i, c in enumerate(table[0])]
                    for row_data in table[1:]:
                        row = {}
                        for i, val in enumerate(row_data):
                            col_name = columns[i] if i < len(columns) else f"col_{i}"
                            row[col_name] = (val or "").strip()
                        if any(v for v in row.values()):
                            rows.append(row)
                else:
                    text = page.extract_text()
                    if text:
                        lines = [l.strip() for l in text.split("\n") if l.strip()]
                        for line in lines:
                            parts = re.split(r"\s{2,}|\t", line)
                            if len(parts) >= 2:
                                if not columns:
                                    columns = [f"col_{i}" for i in range(len(parts))]
                                row = {columns[i] if i < len(columns) else f"col_{i}": parts[i] for i in range(len(parts))}
                                rows.append(row)

        return columns, rows

    def _parse_ppt(self, file_path: str) -> tuple[list[str], list[dict]]:
        from pptx import Presentation
        rows: list[dict] = []
        columns: list[str] = []

        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    if not columns:
                        columns = [cell.text.strip() if cell.text.strip() else f"col_{i}" for i, cell in enumerate(table.rows[0].cells)]
                    for row_idx in range(1, len(table.rows)):
                        row = {}
                        for col_idx, cell in enumerate(table.rows[row_idx].cells):
                            col_name = columns[col_idx] if col_idx < len(columns) else f"col_{col_idx}"
                            row[col_name] = cell.text.strip()
                        if any(v for v in row.values()):
                            rows.append(row)
                elif shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text and "\t" in text:
                        lines = text.strip().split("\n")
                        for line in lines:
                            parts = line.split("\t")
                            if len(parts) >= 2:
                                if not columns:
                                    columns = [f"col_{i}" for i in range(len(parts))]
                                row = {columns[i] if i < len(columns) else f"col_{i}": parts[i] for i in range(len(parts))}
                                rows.append(row)

        return columns, rows

    def _parse_docx(self, file_path: str) -> tuple[list[str], list[dict]]:
        from docx import Document
        rows: list[dict] = []
        columns: list[str] = []

        doc = Document(file_path)
        for table in doc.tables:
            if not columns:
                columns = [cell.text.strip() if cell.text.strip() else f"col_{i}" for i, cell in enumerate(table.rows[0].cells)]
            for row_idx in range(1, len(table.rows)):
                row = {}
                for col_idx, cell in enumerate(table.rows[row_idx].cells):
                    col_name = columns[col_idx] if col_idx < len(columns) else f"col_{col_idx}"
                    row[col_name] = cell.text.strip()
                if any(v for v in row.values()):
                    rows.append(row)

        return columns, rows

    def load_file_data(self, file_name: str) -> list[dict]:
        file_path = self.upload_dir / file_name
        _, data = self.parse_file(str(file_path))
        return data

    def merge_data(
        self,
        rows: list[dict],
        crossref_rows: list[dict],
        match_keys: list[Any],
        output_columns: list[str],
    ) -> list[dict]:
        # Normalizar match_keys a lista de dicts
        keys_list = []
        for m in match_keys:
            if hasattr(m, "model_dump"):
                keys_list.append(m.model_dump())
            elif hasattr(m, "dict"):
                keys_list.append(m.dict())
            else:
                keys_list.append(m)

        ext_keys = [m["extractionKey"] for m in keys_list]
        ref_keys = [m["crossrefKey"] for m in keys_list]

        lookup: dict[tuple, dict] = {}
        for cr_row in crossref_rows:
            key = tuple(str(cr_row.get(k, "")).strip().lower() for k in ref_keys)
            if any(key):
                lookup[key] = cr_row

        merged: list[dict] = []
        for row in rows:
            new_row = dict(row)
            row_key = tuple(str(row.get(k, "")).strip().lower() for k in ext_keys)
            matched = lookup.get(row_key)
            if matched:
                for col in output_columns:
                    new_row[col] = matched.get(col, "")
            else:
                for col in output_columns:
                    new_row[col] = "NO ENCONTRADO"
            merged.append(new_row)

        return merged
